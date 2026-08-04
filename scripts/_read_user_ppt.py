# -*- coding: utf-8 -*-
import sys, glob
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

target = r"C:\Users\youm\Downloads\LLM-AE-AI-mid-ppt.pptx"
candidates = glob.glob(r"C:\Users\youm\Downloads\*.pptx")
for c in candidates:
    if "mid-ppt" in c:
        target = c
        break

print(f"FILE: {target}")
p = Presentation(target)
print(f"#slides: {len(p.slides)}")
print()

for si, s in enumerate(p.slides, 1):
    print(f"========== Slide {si} ==========")
    for sh in s.shapes:
        if sh.has_text_frame:
            tf = sh.text_frame
            for para in tf.paragraphs:
                line = "".join(r.text for r in para.runs).strip()
                if line:
                    print(f"  • {line}")
        elif sh.shape_type == 19:  # TABLE
            try:
                tbl = sh.table
                print("  [TABLE]")
                for row in tbl.rows:
                    cells = [c.text.strip().replace("\n", " | ") for c in row.cells]
                    print("    " + " || ".join(cells))
            except Exception:
                pass
    print()
