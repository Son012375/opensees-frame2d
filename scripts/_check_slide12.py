# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

p = Presentation(r"C:\Users\youm\Desktop\백장운_수업\LLM-AE-AI-mid-ppt_v2.pptx")
print(f"#slides: {len(p.slides)}")
s = p.slides[11]
print(f"Slide 12 — layout: {s.slide_layout.name}, shapes: {len(s.shapes)}")
for sh in s.shapes:
    if sh.has_text_frame:
        txt = " | ".join(p.text.strip() for p in sh.text_frame.paragraphs if p.text.strip())
        if txt:
            print(f"  • {txt}")
