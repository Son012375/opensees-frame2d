# -*- coding: utf-8 -*-
"""기존 LLM-AE-AI-mid-ppt에 백업 슬라이드 12 추가"""
import sys, shutil
from pathlib import Path
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu, Pt, Cm, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SRC = Path(r"C:\Users\youm\Desktop\백장운_수업\LLM-AE-AI-mid-ppt (1).pptx")
DST = Path(r"C:\Users\youm\Desktop\백장운_수업\LLM-AE-AI-mid-ppt_v2.pptx")

NAVY = RGBColor(0x11, 0x2D, 0x4E)
BLUE = RGBColor(0x3F, 0x72, 0xAF)
WHITE = RGBColor(0xF9, 0xF7, 0xF7)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
FONT = "맑은 고딕"


def set_text(tf, text, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # 첫 paragraph 사용
    if not tf.paragraphs:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
        p.clear()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    # East Asian font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", FONT)


def add_text_lines(tf, lines, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_size=None):
    """여러 줄 텍스트. lines = [(text, size, bold, color), ...] 또는 [str, ...]"""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # 첫 줄
    first = lines[0]
    if isinstance(first, str):
        first_data = (first, size, bold, color)
    else:
        first_data = first
    if not tf.paragraphs:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
        p.clear()
    _set_run(p, first_data, align)
    # 나머지 줄
    for line in lines[1:]:
        if isinstance(line, str):
            data = (line, line_size or size, bold, color)
        else:
            data = line
        p2 = tf.add_paragraph()
        _set_run(p2, data, align)


def _set_run(p, data, align):
    txt, size, bold, color = data
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", FONT)


def add_box(slide, left_cm, top_cm, w_cm, h_cm, fill=None, line=None, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, Cm(left_cm), Cm(top_cm), Cm(w_cm), Cm(h_cm))
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_textbox(slide, left_cm, top_cm, w_cm, h_cm):
    tb = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(w_cm), Cm(h_cm))
    return tb


def add_line(slide, x1_cm, y1_cm, x2_cm, y2_cm, color=NAVY, width=1.0):
    line = slide.shapes.add_connector(1, Cm(x1_cm), Cm(y1_cm), Cm(x2_cm), Cm(y2_cm))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def main():
    shutil.copyfile(SRC, DST)
    prs = Presentation(str(DST))
    # 슬라이드 크기 (13.33 x 7.50 in = 33.87 x 19.05 cm)
    print(f"slide size: {Emu(prs.slide_width).cm:.2f} x {Emu(prs.slide_height).cm:.2f} cm")

    # blank layout 찾기 (제목·내용 레이아웃과 동일하게)
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    # 빈 레이아웃 시도
    blank = None
    for lo in prs.slide_layouts:
        if lo.name in ("빈 화면", "Blank"):
            blank = lo
            break
    if blank is None:
        blank = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)

    # placeholder 제거 (있다면)
    for shp in list(slide.placeholders):
        sp = shp._element
        sp.getparent().remove(sp)

    SW = 33.87  # slide width cm
    SH = 19.05  # slide height cm

    # ───── 헤더 (학회 PPT 스타일 매칭) ─────
    # 좌상단 직사각형 액센트 (다크네이비 얇은 막대)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.0), Cm(1.0), Cm(0.5), Cm(1.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = NAVY
    accent.line.fill.background()
    accent.shadow.inherit = False

    # 섹션명 "Appendix"
    tb = add_textbox(slide, 1.8, 0.9, 8.0, 1.0)
    set_text(tb.text_frame, "Appendix", size=28, bold=True, color=NAVY, anchor=MSO_ANCHOR.TOP)
    # 섹션 번호 "A1"
    tb = add_textbox(slide, 1.8, 1.7, 4.0, 1.4)
    set_text(tb.text_frame, "A1", size=44, bold=True, color=BLUE, anchor=MSO_ANCHOR.TOP)
    # 부제
    tb = add_textbox(slide, 5.0, 2.2, 18.0, 0.9)
    set_text(tb.text_frame, "성능 평가 메트릭 (Backup)", size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.TOP)

    # 헤더 하단 구분선
    add_line(slide, 1.0, 3.5, SW - 1.0, 3.5, color=BLUE, width=1.0)

    # ───── 4개 KPI 박스 (F1~F4) ─────
    box_top = 4.0
    box_h = 6.2
    box_w = 7.4
    gap = 0.3
    box_left_start = (SW - (box_w * 4 + gap * 3)) / 2  # 중앙 정렬

    boxes = [
        (
            "F1  KDS-RAG 챗봇",
            [
                ("Top-5 Accuracy", "≥ 0.95"),
                ("MRR", "≥ 0.80"),
                ("Latency (cached)", "≤ 3 s"),
            ],
        ),
        (
            "F2  Design Check 인용",
            [
                ("Hallucination Rate", "= 0"),
                ("Citation Precision", "≥ 0.90"),
                ("Citation Recall", "≥ 0.85"),
            ],
        ),
        (
            "F3  설계 제안 에이전트",
            [
                ("Convergence Rate", "≥ 80 %"),
                ("Time-to-OK", "≤ 60 s"),
                ("Avg Iterations", "≤ 2.5"),
            ],
        ),
        (
            "F4  LLM-as-Judge",
            [
                ("False Positive Rate", "≤ 0.10"),
                ("Cohen's Kappa", "≥ 0.70"),
                ("Self-Consistency", "≥ 0.90"),
            ],
        ),
    ]

    for i, (title, kpis) in enumerate(boxes):
        left = box_left_start + i * (box_w + gap)
        # 박스 (흰 채우기, 파란 테두리)
        box = add_box(slide, left, box_top, box_w, box_h, fill=WHITE, line=BLUE, line_w=1.5)
        # 박스 제목 영역 (상단 0.9cm 다크네이비)
        title_band = add_box(slide, left, box_top, box_w, 1.1, fill=NAVY, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # 다크네이비 띠 위에 흰 글씨로 제목
        tb = add_textbox(slide, left + 0.2, box_top + 0.1, box_w - 0.4, 0.9)
        set_text(tb.text_frame, title, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # KPI 3개 — 각 1.5cm 높이로 배치
        kpi_top = box_top + 1.6
        for j, (name, val) in enumerate(kpis):
            y = kpi_top + j * 1.5
            # KPI 이름
            tb = add_textbox(slide, left + 0.3, y, box_w - 0.6, 0.6)
            set_text(tb.text_frame, name, size=11, bold=False, color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
            # KPI 값 (강조)
            tb = add_textbox(slide, left + 0.3, y + 0.55, box_w - 0.6, 0.7)
            set_text(tb.text_frame, val, size=15, bold=True, color=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # ───── E2E 통합 가로 박스 ─────
    e2e_top = box_top + box_h + 0.4
    e2e_h = 1.6
    e2e_left = 1.0
    e2e_w = SW - 2.0
    e2e_box = add_box(slide, e2e_left, e2e_top, e2e_w, e2e_h, fill=NAVY, line=None)
    # 좌측 라벨 "E2E 통합"
    tb = add_textbox(slide, e2e_left + 0.4, e2e_top, 4.5, e2e_h)
    set_text(tb.text_frame, "E2E 통합", size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # 우측 KPI 4개 가로 배열
    e2e_kpis = [
        ("IFC-to-OK", "≤ 5 min"),
        ("Cost / Run", "≤ $0.5"),
        ("Cache Hit", "≥ 80 %"),
        ("Click Reduction", "≥ 70 %"),
    ]
    kpi_w = (e2e_w - 5.0) / 4
    for i, (name, val) in enumerate(e2e_kpis):
        left = e2e_left + 5.0 + i * kpi_w
        tb = add_textbox(slide, left, e2e_top + 0.15, kpi_w - 0.2, 0.6)
        set_text(tb.text_frame, name, size=10, bold=False, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        tb = add_textbox(slide, left, e2e_top + 0.75, kpi_w - 0.2, 0.8)
        set_text(tb.text_frame, val, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # ───── 평가 셋 박스 ─────
    eval_top = e2e_top + e2e_h + 0.4
    eval_h = 1.4
    eval_box = add_box(slide, 1.0, eval_top, 22.0, eval_h, fill=LIGHT_GRAY, line=None, shape=MSO_SHAPE.RECTANGLE)
    tb = add_textbox(slide, 1.3, eval_top + 0.1, 21.4, eval_h - 0.2)
    add_text_lines(
        tb.text_frame,
        [
            ("[평가 셋]   KDS Gold Set 50~100  ·  NG Case Set 10~20  ·  Human Baseline 5~10   (W2~W4 구축)", 11, True, NAVY),
        ],
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    # ───── 하단 takeaway 메시지 ─────
    msg_top = eval_top + 0.05
    tb = add_textbox(slide, 1.0, SH - 1.6, SW - 2.0, 1.0)
    set_text(
        tb.text_frame,
        "기능별 정량 측정  ·  Hallucination 0% 시스템 차단  ·  Human Baseline 대비 검증",
        size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # 슬라이드 번호 (우하단)
    tb = add_textbox(slide, SW - 2.5, SH - 0.9, 2.0, 0.6)
    set_text(tb.text_frame, "12", size=12, color=GRAY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)

    prs.save(str(DST))
    print(f"[OK] saved: {DST}")
    print(f"slides total: {len(prs.slides)}")


if __name__ == "__main__":
    main()
