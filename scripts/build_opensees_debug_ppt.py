"""
opensees-debug 스킬 소개 PPTX 생성 스크립트
출력: outputs/opensees_debug_skill.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 ──────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x2A, 0x4A)   # 딥 네이비
C_ACCENT = RGBColor(0xE8, 0x6A, 0x1A)   # 오렌지
C_MID    = RGBColor(0x2E, 0x4D, 0x7B)   # 미드 블루
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xFA)   # 연한 블루-화이트
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_GREEN  = RGBColor(0x27, 0x8F, 0x5F)
C_CODE   = RGBColor(0xF5, 0xF5, 0xF5)
C_CODE_T = RGBColor(0x1A, 0x2A, 0x4A)

W  = Inches(13.33)   # 와이드 16:9
H  = Inches(7.5)


# ── 헬퍼 함수들 ─────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.line.color.rgb = fill
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
             font_name="Malgun Gothic", wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, space_before=0, indent=0,
             font_name="Malgun Gothic"):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = _Pt(space_before)
    p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = _Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def header_band(slide, title, subtitle=None):
    """상단 헤더 밴드"""
    add_rect(slide, 0, 0, W, Inches(1.4), C_DARK)
    add_text(slide, title,
             Inches(0.5), Inches(0.18), Inches(11), Inches(0.75),
             font_size=28, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.92), Inches(11), Inches(0.4),
                 font_size=14, color=RGBColor(0xAA, 0xBB, 0xDD), font_name="Malgun Gothic")
    # 오렌지 하단 선
    add_rect(slide, 0, Inches(1.4), W, Pt(4), C_ACCENT)


def slide_num(slide, num, total=7):
    add_text(slide, f"{num} / {total}",
             Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
             font_size=11, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ── 슬라이드 제작 ────────────────────────────────────────

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # 완전 빈 레이아웃


# ════════════════════════════════════════════════════════
# Slide 1 — Title
# ════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1, C_DARK)

# 배경 장식 — 우측 삼각형 느낌 직사각형
add_rect(s1, Inches(9.5), 0, Inches(3.83), H, C_MID)
add_rect(s1, Inches(12.3), 0, Inches(1.03), H, C_ACCENT)

# 스킬 이름
add_text(s1, "opensees-debug",
         Inches(0.7), Inches(1.8), Inches(8.5), Inches(1.4),
         font_size=52, bold=True, color=C_WHITE, font_name="Consolas")

# 부제목
add_text(s1, "OpenSees 구조해석 오류 자동 진단 Claude Code 스킬",
         Inches(0.7), Inches(3.2), Inches(8.5), Inches(0.7),
         font_size=20, color=RGBColor(0xAA, 0xCC, 0xFF), font_name="Malgun Gothic")

# 태그들
tags = ["Singular Matrix", "Convergence Failure", "NaN Detection", "3D Release DOF"]
for i, tag in enumerate(tags):
    tx = Inches(0.7) + i * Inches(2.1)
    add_rect(s1, tx, Inches(4.1), Inches(1.9), Inches(0.42), C_ACCENT)
    add_text(s1, tag, tx + Inches(0.08), Inches(4.12), Inches(1.8), Inches(0.4),
             font_size=11, bold=True, color=C_WHITE, font_name="Malgun Gothic")

# 구분선
add_rect(s1, Inches(0.7), Inches(4.75), Inches(7.5), Pt(2), C_ACCENT)

# 메타
add_text(s1, "Week 05 Agent Skills 과제  |  백정원  |  2026",
         Inches(0.7), Inches(5.0), Inches(8), Inches(0.4),
         font_size=13, color=RGBColor(0x88, 0xAA, 0xCC), font_name="Malgun Gothic")

# 우측 패널 텍스트
add_text(s1, "Claude Code\nAgent Skill",
         Inches(9.8), Inches(3.0), Inches(3.0), Inches(1.2),
         font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")
add_text(s1, "18 Error Patterns\n5-Step Diagnosis\nAuto Health Check",
         Inches(9.8), Inches(4.4), Inches(3.0), Inches(1.5),
         font_size=14, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER, font_name="Malgun Gothic")


# ════════════════════════════════════════════════════════
# Slide 2 — 문제 상황
# ════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
set_bg(s2, C_WHITE)
header_band(s2, "문제 상황", "OpenSees 오류, 왜 이렇게 디버깅이 어려울까?")
slide_num(s2, 2)

# 왼쪽 — 오류 메시지 박스
add_rect(s2, Inches(0.4), Inches(1.65), Inches(6.2), Inches(4.8), RGBColor(0x1E, 0x1E, 0x1E))
add_text(s2, "$ python run_analysis.py",
         Inches(0.55), Inches(1.75), Inches(6.0), Inches(0.35),
         font_size=11, color=RGBColor(0x88, 0xFF, 0x88), font_name="Consolas")

error_lines = [
    ("WARNING: SparseGenRowLinSOE::", RGBColor(0xFF, 0x77, 0x77)),
    ("  solve() - singular matrix detected", RGBColor(0xFF, 0x77, 0x77)),
    ("", C_WHITE),
    ("StaticAnalysis::analyze() failed", RGBColor(0xFF, 0xAA, 0x44)),
    ("  - the algorithm failed to converge", RGBColor(0xFF, 0xAA, 0x44)),
    ("  after 50 iterations", RGBColor(0xFF, 0xAA, 0x44)),
    ("", C_WHITE),
    ("ok = -1", RGBColor(0xFF, 0x55, 0x55)),
]
for idx, (line, col) in enumerate(error_lines):
    add_text(s2, line,
             Inches(0.6), Inches(2.15) + idx * Inches(0.38), Inches(5.8), Inches(0.4),
             font_size=11, color=col, font_name="Consolas")

# 오류 박스 레이블
add_rect(s2, Inches(0.4), Inches(1.65), Inches(2.2), Inches(0.32), RGBColor(0x44, 0x44, 0x44))
add_text(s2, "  터미널 오류 출력",
         Inches(0.42), Inches(1.66), Inches(2.1), Inches(0.3),
         font_size=10, color=RGBColor(0xCC, 0xCC, 0xCC), font_name="Malgun Gothic")

# 오른쪽 — 문제점 설명
add_rect(s2, Inches(7.0), Inches(1.65), Inches(5.9), Inches(4.8), C_LIGHT)
add_text(s2, "개발자가 겪는 현실",
         Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.45),
         font_size=16, bold=True, color=C_DARK, font_name="Malgun Gothic")

problems = [
    ("오류 메시지 ≠ 실제 원인",   "singular matrix는 결과, 진짜 원인은\n경계조건·릴리즈·부유노드 등 여러 가지"),
    ("검색해도 맥락 없는 답변",   "일반 OpenSees 포럼은 영어,\n이 프로젝트 구조를 모름"),
    ("반복적인 시행착오",          "경계조건 추가 → 릴리즈 수정 → 하중 재확인\n→ 30분 이상 소요"),
]
for idx, (title, desc) in enumerate(problems):
    y = Inches(2.35) + idx * Inches(1.3)
    add_rect(s2, Inches(7.15), y, Pt(5), Inches(0.9), C_ACCENT)
    add_text(s2, title, Inches(7.35), y, Inches(5.3), Inches(0.38),
             font_size=13, bold=True, color=C_DARK, font_name="Malgun Gothic")
    add_text(s2, desc, Inches(7.35), y + Inches(0.37), Inches(5.3), Inches(0.65),
             font_size=11, color=C_GRAY, font_name="Malgun Gothic")


# ════════════════════════════════════════════════════════
# Slide 3 — 스킬 소개
# ════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
set_bg(s3, C_WHITE)
header_band(s3, "opensees-debug 스킬 소개", "Claude에게 OpenSees 전문 지식을 주입")
slide_num(s3, 3)

# 중앙 큰 설명 카드 3개
cards = [
    (C_DARK,   "WHAT",
     "OpenSees 해석 오류를 자동으로 진단하는\nClaude Code 스킬",
     "18가지 오류 패턴 인식, 5단계 진단 절차,\n자동 진단 스크립트 포함"),
    (C_MID,    "WHEN",
     "자동 트리거 — 키워드 감지 시 즉시 활성화",
     '"singular matrix 왜 나지?"\n"OpenSees가 수렴을 못해"\n"변위가 너무 크게 나온다"'),
    (C_ACCENT, "HOW",
     "3레이어 구조로 컨텍스트 절약",
     "Level 1: name+description (항상 로드)\nLevel 2: SKILL.md (매칭 시)\nLevel 3: references/ (필요 시)"),
]
for i, (col, label, title, body) in enumerate(cards):
    x = Inches(0.35) + i * Inches(4.32)
    add_rect(s3, x, Inches(1.65), Inches(4.1), Inches(5.1), col)
    # 레이블
    add_text(s3, label, x + Inches(0.15), Inches(1.75), Inches(3.8), Inches(0.55),
             font_size=22, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    # 구분선
    add_rect(s3, x + Inches(0.15), Inches(2.32), Inches(3.8), Pt(2),
             RGBColor(0xFF, 0xFF, 0xFF) if col != C_ACCENT else C_DARK)
    # 제목
    add_text(s3, title, x + Inches(0.15), Inches(2.45), Inches(3.8), Inches(0.9),
             font_size=14, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    # 본문
    add_text(s3, body, x + Inches(0.15), Inches(3.4), Inches(3.8), Inches(3.0),
             font_size=12, color=RGBColor(0xDD, 0xEE, 0xFF) if col != C_ACCENT else C_WHITE,
             font_name="Malgun Gothic")


# ════════════════════════════════════════════════════════
# Slide 4 — 스킬 파일 구조
# ════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
set_bg(s4, C_WHITE)
header_band(s4, "스킬 파일 구조", "Progressive Disclosure — 필요할 때만 로드")
slide_num(s4, 4)

# 왼쪽 — 디렉토리 트리
add_rect(s4, Inches(0.4), Inches(1.65), Inches(5.8), Inches(5.6), RGBColor(0x1E, 0x1E, 0x1E))
tree_lines = [
    ("  .claude/skills/opensees-debug/", RGBColor(0xFF, 0xDD, 0x88)),
    ("  ├── SKILL.md          ← Level 2", RGBColor(0x88, 0xFF, 0x88)),
    ("  │    (오류 분류표 + 5단계 절차)", RGBColor(0x88, 0xAA, 0x88)),
    ("  ├── references/       ← Level 3", RGBColor(0x88, 0xCC, 0xFF)),
    ("  │   ├── error_catalog.md", RGBColor(0x88, 0xCC, 0xFF)),
    ("  │   │    (18가지 오류 패턴)", RGBColor(0x66, 0x99, 0xBB)),
    ("  │   ├── model_checklist.md", RGBColor(0x88, 0xCC, 0xFF)),
    ("  │   └── fix_patterns.md", RGBColor(0x88, 0xCC, 0xFF)),
    ("  └── scripts/          ← Level 3", RGBColor(0xFF, 0xAA, 0x88)),
    ("      └── check_model_health.py", RGBColor(0xFF, 0xAA, 0x88)),
    ("           (자동 진단 스크립트)", RGBColor(0xCC, 0x88, 0x66)),
]
for idx, (line, col) in enumerate(tree_lines):
    add_text(s4, line,
             Inches(0.5), Inches(1.85) + idx * Inches(0.43), Inches(5.6), Inches(0.42),
             font_size=11, color=col, font_name="Consolas")

# 오른쪽 — Progressive Disclosure 설명
levels = [
    ("Level 1", "항상 로드", "name + description만\n(수십 토큰)",
     "스킬 매칭 판단용", C_GREEN),
    ("Level 2", "매칭 시 로드", "SKILL.md 본문\n(수백 토큰)",
     "오류 분류 + 진단 절차", C_MID),
    ("Level 3", "필요 시 로드", "references/ + scripts/\n(수천 토큰)",
     "상세 카탈로그 + 코드", C_ACCENT),
]
for i, (lv, when, content, purpose, col) in enumerate(levels):
    y = Inches(1.7) + i * Inches(1.75)
    add_rect(s4, Inches(6.5), y, Inches(0.75), Inches(1.5), col)
    add_text(s4, lv, Inches(6.5), y + Inches(0.45), Inches(0.75), Inches(0.55),
             font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")
    add_rect(s4, Inches(7.3), y, Inches(5.55), Inches(1.5), C_LIGHT)
    add_text(s4, when, Inches(7.45), y + Inches(0.08), Inches(5.2), Inches(0.4),
             font_size=13, bold=True, color=col, font_name="Malgun Gothic")
    add_text(s4, content, Inches(7.45), y + Inches(0.48), Inches(2.6), Inches(0.7),
             font_size=11, color=C_DARK, font_name="Consolas")
    add_text(s4, purpose, Inches(10.1), y + Inches(0.48), Inches(2.6), Inches(0.7),
             font_size=11, color=C_GRAY, font_name="Malgun Gothic")

    if i < 2:
        add_rect(s4, Inches(6.85), y + Inches(1.5), Pt(5), Inches(0.25), col)


# ════════════════════════════════════════════════════════
# Slide 5 — 5단계 진단 프로세스
# ════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
set_bg(s5, C_WHITE)
header_band(s5, "5단계 진단 프로세스", "오류 발생 → 자동 진단 → 수정 완료")
slide_num(s5, 5)

steps = [
    ("1", "오류 메시지 분류",
     "키워드 패턴 매칭\n→ E-01 ~ E-18 코드 부여",
     "Grep, Read"),
    ("2", "모델 자동 진단",
     "check_model_health.py 실행\n소스코드/결과JSON 정적 분석",
     "Bash"),
    ("3", "원인 확정",
     "model_checklist.md 순서대로\n5가지 필수 항목 체크",
     "Read"),
    ("4", "수정 코드 제시",
     "fix_patterns.md에서\n해당 유형 패턴 적용",
     "Read, Edit"),
    ("5", "수정 후 검증",
     "ok==0 확인, 반력 평형,\nNaN 탐지 자동 실행",
     "Bash"),
]

arrow_col = RGBColor(0xCC, 0xCC, 0xCC)
for i, (num, title, body, tools) in enumerate(steps):
    x = Inches(0.25) + i * Inches(2.58)
    # 박스
    col = C_ACCENT if i == 0 else (C_MID if i == 4 else C_DARK)
    add_rect(s5, x, Inches(1.65), Inches(2.4), Inches(0.65), col)
    add_text(s5, f"Step {num}  {title}",
             x + Inches(0.1), Inches(1.7), Inches(2.3), Inches(0.55),
             font_size=12, bold=True, color=C_WHITE, font_name="Malgun Gothic")
    # 화살표
    if i < 4:
        add_rect(s5, x + Inches(2.4), Inches(1.9), Inches(0.18), Pt(3), arrow_col)
        add_text(s5, "▶", x + Inches(2.4), Inches(1.82), Inches(0.18), Inches(0.35),
                 font_size=12, color=arrow_col, align=PP_ALIGN.CENTER, font_name="Malgun Gothic")
    # 본문
    add_rect(s5, x, Inches(2.35), Inches(2.4), Inches(3.8), C_LIGHT)
    add_text(s5, body,
             x + Inches(0.1), Inches(2.45), Inches(2.25), Inches(2.3),
             font_size=11, color=C_DARK, font_name="Malgun Gothic")
    # 툴 뱃지
    add_rect(s5, x + Inches(0.05), Inches(5.7), Inches(2.3), Inches(0.32), col)
    add_text(s5, f"tools: {tools}",
             x + Inches(0.1), Inches(5.73), Inches(2.25), Inches(0.3),
             font_size=10, color=C_WHITE, font_name="Consolas")

# 하단 주석
add_text(s5, "* allowed-tools: Read, Grep, Bash  (SKILL.md frontmatter 설정)",
         Inches(0.3), Inches(6.9), Inches(12), Inches(0.35),
         font_size=11, color=C_GRAY, font_name="Malgun Gothic")


# ════════════════════════════════════════════════════════
# Slide 6 — 데모: Singular Matrix 해결
# ════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
set_bg(s6, C_WHITE)
header_band(s6, "데모: Singular Matrix 오류 해결", "경계조건 누락 → E-01 진단 → fix() 추가 완료")
slide_num(s6, 6)

# BEFORE
add_rect(s6, Inches(0.35), Inches(1.65), Inches(6.0), Inches(0.4), C_RED)
add_text(s6, "  BEFORE  오류 발생",
         Inches(0.35), Inches(1.66), Inches(5.9), Inches(0.38),
         font_size=13, bold=True, color=C_WHITE, font_name="Malgun Gothic")
add_rect(s6, Inches(0.35), Inches(2.05), Inches(6.0), Inches(2.3), RGBColor(0x1E, 0x1E, 0x1E))
before_lines = [
    ("ops.wipe()",                         C_WHITE),
    ("ops.model('basic', '-ndm', 3, '-ndf', 6)", C_WHITE),
    ("# 노드 정의...",                     RGBColor(0x77, 0x99, 0x77)),
    ("# 요소 정의...",                     RGBColor(0x77, 0x99, 0x77)),
    ("# ← ops.fix() 없음!",               RGBColor(0xFF, 0x55, 0x55)),
    ("ops.analyze(1)  # → ok = -1",        RGBColor(0xFF, 0x77, 0x44)),
]
for idx, (line, col) in enumerate(before_lines):
    add_text(s6, f"  {line}",
             Inches(0.4), Inches(2.1) + idx * Inches(0.34), Inches(5.9), Inches(0.35),
             font_size=10.5, color=col, font_name="Consolas")

# AFTER
add_rect(s6, Inches(7.0), Inches(1.65), Inches(6.0), Inches(0.4), C_GREEN)
add_text(s6, "  AFTER  수정 완료",
         Inches(7.0), Inches(1.66), Inches(5.9), Inches(0.38),
         font_size=13, bold=True, color=C_WHITE, font_name="Malgun Gothic")
add_rect(s6, Inches(7.0), Inches(2.05), Inches(6.0), Inches(2.3), RGBColor(0x1E, 0x1E, 0x1E))
after_lines = [
    ("ops.wipe()",                          C_WHITE),
    ("ops.model('basic', '-ndm', 3, ...)",  C_WHITE),
    ("# 노드 정의...",                      RGBColor(0x77, 0x99, 0x77)),
    ("# 요소 정의...",                      RGBColor(0x77, 0x99, 0x77)),
    ("for n in base_nodes:",                RGBColor(0x88, 0xCC, 0xFF)),
    ("    ops.fix(n, 1,1,1,1,1,1)  # ✓",   RGBColor(0x88, 0xFF, 0x88)),
    ("ok = ops.analyze(1)  # ok = 0",       RGBColor(0x88, 0xFF, 0x88)),
]
for idx, (line, col) in enumerate(after_lines):
    add_text(s6, f"  {line}",
             Inches(7.05), Inches(2.1) + idx * Inches(0.34), Inches(5.9), Inches(0.35),
             font_size=10.5, color=col, font_name="Consolas")

# 화살표
add_text(s6, "→",
         Inches(6.15), Inches(2.8), Inches(0.8), Inches(0.6),
         font_size=36, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# 하단 진단 흐름
add_rect(s6, Inches(0.35), Inches(4.45), Inches(12.65), Inches(2.7), C_LIGHT)
add_text(s6, "스킬 실행 흐름",
         Inches(0.55), Inches(4.55), Inches(4), Inches(0.4),
         font_size=13, bold=True, color=C_DARK, font_name="Malgun Gothic")

flow_items = [
    ("사용자 입력", '"singular matrix 오류가 왜 나지?"'),
    ("Step 1 매칭", "→ E-01: SparseGenRowLinSOE singular"),
    ("Step 2 진단", "→ check_model_health.py: [CRITICAL] fix() 없음"),
    ("Step 3 확인", "→ checklist A-1: 경계조건 0개 (기준: ≥1)"),
    ("Step 4 수정", "→ fix_patterns.md E-01: for n in base_nodes: ops.fix(n, ...)"),
    ("Step 5 검증", "→ ok=0, 반력합=총하중 ✓"),
]
for i, (label, val) in enumerate(flow_items):
    x = Inches(0.55) + (i % 3) * Inches(4.2)
    y = Inches(5.0) + (i // 3) * Inches(0.55)
    add_text(s6, f"{label}: {val}",
             x, y, Inches(4.1), Inches(0.48),
             font_size=10, color=C_DARK if i % 2 == 0 else C_MID, font_name="Malgun Gothic")


# ════════════════════════════════════════════════════════
# Slide 7 — 확장 가능성 + 마무리
# ════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
set_bg(s7, C_WHITE)
header_band(s7, "확장 가능성 & 마무리", "이 스킬을 시작점으로")
slide_num(s7, 7)

# 왼쪽 — 확장 카드
add_text(s7, "단기 확장",
         Inches(0.4), Inches(1.75), Inches(6.0), Inches(0.4),
         font_size=15, bold=True, color=C_ACCENT, font_name="Malgun Gothic")
short_items = [
    "Midas Gen 오류 카탈로그 추가",
    "SAP2000 입력 오류 패턴 추가",
    "오류 코드별 유사 이슈 GitHub 링크",
    "단면 DB 미등록 경고 자동화",
]
for i, item in enumerate(short_items):
    add_rect(s7, Inches(0.4), Inches(2.2) + i * Inches(0.62), Pt(6), Inches(0.42), C_ACCENT)
    add_text(s7, item,
             Inches(0.65), Inches(2.22) + i * Inches(0.62), Inches(5.7), Inches(0.42),
             font_size=12, color=C_DARK, font_name="Malgun Gothic")

add_text(s7, "중장기 확장",
         Inches(0.4), Inches(4.75), Inches(6.0), Inches(0.4),
         font_size=15, bold=True, color=C_MID, font_name="Malgun Gothic")
long_items = [
    "CI/CD 연동: 커밋 전 자동 진단",
    "해석 결과 → 노션 자동 업로드",
    "수렴 이력 로그 기반 패턴 학습",
]
for i, item in enumerate(long_items):
    add_rect(s7, Inches(0.4), Inches(5.2) + i * Inches(0.55), Pt(6), Inches(0.4), C_MID)
    add_text(s7, item,
             Inches(0.65), Inches(5.22) + i * Inches(0.55), Inches(5.7), Inches(0.4),
             font_size=12, color=C_DARK, font_name="Malgun Gothic")

# 오른쪽 — 요약 카드
add_rect(s7, Inches(7.0), Inches(1.65), Inches(5.9), Inches(5.5), C_DARK)
add_text(s7, "opensees-debug 요약",
         Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.45),
         font_size=16, bold=True, color=C_WHITE, font_name="Malgun Gothic")
add_rect(s7, Inches(7.2), Inches(2.28), Inches(5.3), Pt(2), C_ACCENT)

summary_items = [
    ("스킬 타입",  "Claude Code Agent Skill"),
    ("트리거",     "자동 (오류 키워드 감지)"),
    ("진단 범위",  "18가지 오류 패턴"),
    ("파일 구성",  "SKILL.md + references(3) + scripts(1)"),
    ("allowed-tools", "Read, Grep, Bash"),
    ("모델",       "sonnet"),
    ("설치",       "~/.claude/skills/ 또는 .claude/skills/"),
    ("GitHub",     "Son012375/opensees-debug-skill"),
]
for i, (k, v) in enumerate(summary_items):
    y = Inches(2.45) + i * Inches(0.55)
    add_text(s7, k,
             Inches(7.2), y, Inches(1.9), Inches(0.45),
             font_size=11, bold=True, color=C_ACCENT, font_name="Malgun Gothic")
    add_text(s7, v,
             Inches(9.2), y, Inches(3.5), Inches(0.45),
             font_size=11, color=C_WHITE, font_name="Malgun Gothic")


# ── 출력 ────────────────────────────────────────────────
out_dir = Path("d:/son/opensees-MCP/outputs")
out_dir.mkdir(exist_ok=True)
out_path = out_dir / "opensees_debug_skill_v2.pptx"
prs.save(str(out_path))
print(f"[OK] PPTX 저장: {out_path}")
