# -*- coding: utf-8 -*-
import sys, os, glob
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

candidates = glob.glob(r"C:\Users\youm\Downloads\*.pptx")
target = None
for c in candidates:
    if "콘학" in c or "발표자료-2" in c:
        target = c
        break
if target is None:
    target = max(candidates, key=os.path.getmtime)

print(f"FILE: {target}")
p = Presentation(target)
print(f"slide_size: {Emu(p.slide_width).inches:.2f} x {Emu(p.slide_height).inches:.2f} in")
print(f"#slides: {len(p.slides)}")
print()

for si, s in enumerate(p.slides, 1):
    print(f"=== Slide {si} (layout={s.slide_layout.name}) ===")
    for sh in s.shapes:
        kind = sh.shape_type
        name = sh.name
        info = f"  [{kind}] {name}"
        if sh.has_text_frame:
            tf = sh.text_frame
            txts = []
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        f = run.font
                        size = f.size.pt if f.size else "?"
                        bold = "B" if f.bold else "-"
                        color = ""
                        try:
                            if f.color and f.color.rgb:
                                color = f"#{f.color.rgb}"
                        except Exception:
                            pass
                        txts.append(f"{run.text!r}({size}/{bold}{color})")
            if txts:
                info += " :: " + " | ".join(txts[:6])
        try:
            if sh.fill and sh.fill.type is not None:
                pass
        except Exception:
            pass
        print(info)
    print()
